import streamlit as st
import json
import os
import re
from pathlib import Path

# Try to load dotenv, but don't fail if not available
try:
    from dotenv import load_dotenv

    load_dotenv()
except:
    pass

# Document processing libraries
try:
    import PyPDF2
    from docx import Document
    from pptx import Presentation
except ImportError as e:
    st.error(f"Missing required library: {e}")
    st.info("Run: pip install PyPDF2 python-docx python-pptx")

# Hugging Face
try:
    from huggingface_hub import InferenceClient
except ImportError as e:
    st.error(f"Missing huggingface-hub library: {e}")
    st.info("Run: pip install huggingface-hub")

# ==================== HUGGING FACE SETUP ====================
# Try to get token from environment or session state
HF_API_TOKEN = os.getenv("HUGGINGFACE_API_TOKEN", "")

# List of free, working models for MCQ generation
HF_MODELS = [
    "google/flan-t5-large",  # Good for structured output
    "microsoft/phi-2",  # 2.7B params, good quality
]


def get_hf_client():
    """Initialize Hugging Face client with proper authentication"""
    if not HF_API_TOKEN:
        return None

    try:
        client = InferenceClient(token=HF_API_TOKEN)
        return client
    except Exception as e:
        st.error(f"Failed to initialize client: {str(e)}")
        return None


def test_hf_connection():
    """Test if Hugging Face connection works"""
    if not HF_API_TOKEN:
        return False
    try:
        client = InferenceClient(token=HF_API_TOKEN)
        return True
    except:
        return False


# ==================== DOCUMENT PROCESSING ====================
def extract_text_from_pdf(file):
    """Extract text from PDF file"""
    text = ""
    try:
        pdf_reader = PyPDF2.PdfReader(file)
        for page in pdf_reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    except Exception as e:
        st.error(f"Error reading PDF: {str(e)}")
    return text


def extract_text_from_docx(file):
    """Extract text from DOCX file"""
    text = ""
    try:
        doc = Document(file)
        for paragraph in doc.paragraphs:
            if paragraph.text:
                text += paragraph.text + "\n"
    except Exception as e:
        st.error(f"Error reading DOCX: {str(e)}")
    return text


def extract_text_from_pptx(file):
    """Extract text from PPTX file"""
    text = ""
    try:
        prs = Presentation(file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
    except Exception as e:
        st.error(f"Error reading PPTX: {str(e)}")
    return text


def process_uploaded_file(uploaded_file):
    """Process uploaded file and extract text based on file type"""
    if uploaded_file is None:
        return None

    file_extension = Path(uploaded_file.name).suffix.lower()

    if file_extension == '.pdf':
        return extract_text_from_pdf(uploaded_file)
    elif file_extension == '.docx':
        return extract_text_from_docx(uploaded_file)
    elif file_extension == '.pptx':
        return extract_text_from_pptx(uploaded_file)
    else:
        st.error(f"Unsupported file type: {file_extension}")
        return None


# ==================== FALLBACK QUESTION GENERATOR ====================
def generate_fallback_questions(text_content, quiz_level, num_questions=10):
    """Generate questions using text analysis (no API needed)"""
    import random
    import re

    # Split into sentences
    sentences = re.split(r'[.!?]+', text_content)
    sentences = [s.strip() for s in sentences if len(s.strip()) > 30]

    if len(sentences) < num_questions:
        sentences = sentences * (num_questions // len(sentences) + 1)

    questions = []
    used_indices = set()

    for i in range(min(num_questions, len(sentences))):
        # Get unique sentence
        idx = random.randint(0, len(sentences) - 1)
        while idx in used_indices and len(used_indices) < len(sentences):
            idx = random.randint(0, len(sentences) - 1)
        used_indices.add(idx)

        sentence = sentences[idx]
        words = sentence.split()

        if len(words) > 3:
            # Create question based on sentence
            if len(words) > 5:
                # Remove a word to create fill-in-the-blank
                word_idx = random.randint(0, min(4, len(words) - 1))
                removed_word = words[word_idx]
                question_text = sentence.replace(removed_word, "______", 1)

                # Generate options
                options = {
                    "a": removed_word,
                    "b": words[random.randint(0, len(words) - 1)] if len(words) > 1 else "option",
                    "c": words[random.randint(0, len(words) - 1)] if len(words) > 2 else "option",
                    "d": words[random.randint(0, len(words) - 1)] if len(words) > 3 else "option"
                }

                questions.append({
                    "mcq": f"Complete the sentence: {question_text}",
                    "options": options,
                    "correct": "a"
                })
            else:
                # Comprehension question
                question_text = f"Based on the text, what is mentioned about: '{sentence[:50]}...'?"

                options = {
                    "a": sentence[:80],
                    "b": "Not mentioned in the text",
                    "c": "The opposite is true",
                    "d": "None of the above"
                }

                questions.append({
                    "mcq": question_text,
                    "options": options,
                    "correct": "a"
                })

    return questions[:num_questions]


# ==================== MCQ GENERATION WITH HF ====================
def fetch_questions_with_hf(text_content, quiz_level, num_questions=10, model_index=0):
    """Generate MCQs using Hugging Face Inference API"""

    if not text_content or len(text_content.strip()) < 50:
        return []

    # Limit text length
    if len(text_content) > 2000:
        text_content = text_content[:2000] + "..."

    selected_model = HF_MODELS[model_index] if model_index < len(HF_MODELS) else HF_MODELS[0]

    # Simplified prompt to ensure JSON response
    prompt = f"""Generate {num_questions} multiple choice questions from this text. Difficulty: {quiz_level}

Text: {text_content}

Return ONLY this JSON format:
{{"mcqs": [{{"mcq": "question", "options": {{"a": "opt1", "b": "opt2", "c": "opt3", "d": "opt4"}}, "correct": "a"}}]}}

Generate {num_questions} questions now:"""

    try:
        client = get_hf_client()
        if client is None:
            return None

        response = client.chat_completion(
            model=selected_model,
            messages=[
                {"role": "system", "content": "You are a JSON-only quiz generator. Respond with valid JSON only."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.3,
            max_tokens=3000,
        )

        raw_response = response.choices[0].message.content

        # Clean response
        raw_response = raw_response.strip()
        if raw_response.startswith('```json'):
            raw_response = raw_response[7:]
        if raw_response.startswith('```'):
            raw_response = raw_response[3:]
        if raw_response.endswith('```'):
            raw_response = raw_response[:-3]

        # Find JSON
        json_match = re.search(r'\{.*\}', raw_response, re.DOTALL)
        if json_match:
            raw_response = json_match.group()

        result = json.loads(raw_response)

        if "mcqs" in result and len(result["mcqs"]) > 0:
            return result["mcqs"]
        else:
            return None

    except Exception as e:
        st.warning(f"HF API error: {str(e)[:100]}")
        return None


@st.cache_data
def fetch_questions(text_content, quiz_level, num_questions=10, model_index=0, use_fallback=True):
    """Main function to fetch questions - tries HF first, then fallback"""

    # Try Hugging Face first
    questions = fetch_questions_with_hf(text_content, quiz_level, num_questions, model_index)

    # If HF fails and fallback is enabled, use fallback generator
    if (questions is None or len(questions) == 0) and use_fallback:
        st.info("Using smart text analysis to generate questions...")
        questions = generate_fallback_questions(text_content, quiz_level, num_questions)

    return questions if questions else []


# ==================== EXPORT FUNCTIONS ====================
def export_questions_to_json(questions):
    """Export questions to JSON format"""
    export_data = {
        "quiz_title": "Generated Quiz",
        "total_questions": len(questions),
        "questions": questions
    }
    return json.dumps(export_data, indent=2)


def export_questions_to_txt(questions):
    """Export questions to text format"""
    output = []
    output.append("=" * 60)
    output.append("GENERATED QUIZ QUESTIONS")
    output.append("=" * 60)
    output.append("")

    for idx, q in enumerate(questions, 1):
        output.append(f"Question {idx}: {q['mcq']}")
        output.append("-" * 40)
        for opt_key, opt_value in q['options'].items():
            output.append(f"  {opt_key}. {opt_value}")
        output.append(f"\n✓ Correct Answer: {q['correct'].upper()}")
        output.append("")
        output.append("=" * 60)
        output.append("")

    return "\n".join(output)


# ==================== TOKEN INPUT DIALOG ====================
def show_token_input():
    """Show token input dialog if no token is found"""
    st.sidebar.markdown("---")
    st.sidebar.subheader("🔑 Hugging Face Setup")

    token = st.sidebar.text_input(
        "Enter Hugging Face Token",
        type="password",
        help="Get your token from https://huggingface.co/settings/tokens"
    )

    if token:
        st.session_state.hf_token = token
        os.environ["HUGGINGFACE_API_TOKEN"] = token
        st.sidebar.success("✅ Token saved for this session!")
        st.sidebar.info("Restart the app and add to .env file for permanent storage")

    st.sidebar.markdown("""
    **How to get a token:**
    1. Go to [huggingface.co/settings/tokens](https://huggingface.co/settings/tokens)
    2. Click "New token"
    3. Select "Inference" permission
    4. Copy the token (starts with `hf_`)
    5. Create `.env` file with:""")

# ==================== MAIN APP ====================
def main():
    st.set_page_config(
 page_title="AI Quiz Generator BY Muhammad Rauf Khan",
 page_icon="📚",
 layout="wide"
)

st.title("📚, AI-Based Quiz Generator by Muhammad Rauf Khan")

# Initialize session state
if 'hf_token' not in st.session_state:
 st.session_state.hf_token = os.getenv("HUGGINGFACE_API_TOKEN", "")

if 'quiz_generated' not in st.session_state:
 st.session_state.quiz_generated = False
if 'quiz_submitted' not in st.session_state:
 st.session_state.quiz_submitted = False
if 'questions' not in st.session_state:
 st.session_state.questions = []
if 'user_answers' not in st.session_state:
 st.session_state.user_answers = {}

# Show token input in sidebar
if not st.session_state.hf_token:
 show_token_input()

st.markdown("---")

# ==================== SIDEBAR ====================
with st.sidebar:
 st.header("⚙️ Settings")

 # Model selection (only if token exists)
 if st.session_state.hf_token:
     model_index = st.selectbox(
         "AI Model",
         range(len(HF_MODELS)),
         format_func=lambda x: HF_MODELS[x].split("/")[-1],
         help="Select which AI model to use"
     )
 else:
     model_index = 0
     st.info("💡 Add Hugging Face token for AI-powered questions")

 # Quiz settings
 quiz_level = st.selectbox(
     "Difficulty Level",
     ["Easy", "Medium", "Hard"],
     help="Select difficulty"
 )

 num_questions = st.slider(
     "Number of Questions",
     min_value=5,
     max_value=25,
     value=10,
     step=5
 )

 use_fallback = st.checkbox(
     "Use Smart Text Analysis (no API)",
     value=not bool(st.session_state.hf_token),
     help="Generate questions from text without AI"
 )

 st.markdown("---")
 st.info("💡 **Tips:**\n- Start with 5-10 questions\n- Use fallback mode for offline use")

# ==================== MAIN CONTENT ====================
col1, col2 = st.columns([1, 1])

with col1:
 st.subheader("📄 Upload Document")

 uploaded_file = st.file_uploader(
     "Choose a file (PDF, DOCX, PPTX)",
     type=['pdf', 'docx', 'pptx']
 )

 st.markdown("**— OR —**")

 text_content = st.text_area(
     "📝 Paste Text",
     height=200,
     placeholder="Paste your content here..."
 )

 if uploaded_file:
     with st.spinner(f"Processing..."):
         extracted = process_uploaded_file(uploaded_file)
         if extracted:
             text_content = extracted
             st.success(f"✅ Extracted {len(extracted)} characters")
             with st.expander("Preview"):
                 st.text(extracted[:500] + ("..." if len(extracted) > 500 else ""))

with col2:
 st.subheader("🎯 Generate Quiz")

 if st.button("🚀 Generate Quiz", type="primary", use_container_width=True):
     if not text_content or len(text_content.strip()) < 50:
         st.error("❌ Please provide at least 50 characters of text!")
     else:
         with st.spinner(f"Generating {num_questions} questions..."):
             questions = fetch_questions(
                 text_content=text_content,
                 quiz_level=quiz_level.lower(),
                 num_questions=num_questions,
                 model_index=model_index if 'model_index' in locals() else 0,
                 use_fallback=use_fallback or not st.session_state.hf_token
             )

         if questions and len(questions) > 0:
             st.session_state.questions = questions
             st.session_state.quiz_generated = True
             st.session_state.quiz_submitted = False
             st.success(f"✅ Generated {len(questions)} questions!")

             # Export
             st.subheader("📥 Export")
             c1, c2 = st.columns(2)
             with c1:
                 st.download_button(
                     "JSON",
                     export_questions_to_json(questions),
                     "quiz.json",
                     "application/json"
                 )
             with c2:
                 st.download_button(
                     "TXT",
                     export_questions_to_txt(questions),
                     "quiz.txt",
                     "text/plain"
                 )
         else:
             st.error("Failed to generate questions. Try fewer questions or use fallback mode.")

# ==================== QUIZ DISPLAY ====================
if st.session_state.quiz_generated and st.session_state.questions:
 st.markdown("---")
 st.header("📋 Take the Quiz")

 questions = st.session_state.questions

 with st.form(key="quiz_form"):
     for idx, q in enumerate(questions):
         st.subheader(f"Q{idx + 1}: {q['mcq']}")
         options = list(q["options"].values())

         selected = st.radio(
             "Answer:",
             options,
             key=f"q_{idx}",
             index=None,
             horizontal=True,
             label_visibility="collapsed"
         )

         if selected:
             for key, val in q["options"].items():
                 if val == selected:
                     st.session_state.user_answers[idx] = key
                     break
         else:
             st.session_state.user_answers[idx] = None

         st.divider()

     submitted = st.form_submit_button("✅ Submit Quiz", type="primary")

 if submitted:
     marks = 0
     for idx, q in enumerate(questions):
         user = st.session_state.user_answers.get(idx)
         if user and user == q["correct"]:
             marks += 1

     percentage = (marks / len(questions)) * 100

     st.balloons()
     st.markdown("---")
     st.header("📊 Results")

     col1, col2, col3 = st.columns(3)
     col1.metric("Total", len(questions))
     col2.metric("Correct", marks)
     col3.metric("Score", f"{percentage:.1f}%")

     st.progress(percentage / 100)

     if percentage >= 80:
         st.success("🎉 Excellent!")
     elif percentage >= 60:
         st.success("👍 Good job!")
     else:
         st.warning("📚 Keep practicing!")

     st.subheader("📝 Details")
     for idx, q in enumerate(questions):
         user_ans = st.session_state.user_answers.get(idx)
         is_correct = user_ans == q["correct"]

         with st.expander(f"Q{idx + 1}: {q['mcq'][:100]}..."):
             st.write(f"**Your answer:** {q['options'].get(user_ans, 'Not answered')}")
             st.write(f"**Correct answer:** {q['options'][q['correct']]}")
             if is_correct:
                 st.success("✅ Correct")
             else:
                 st.error("❌ Incorrect")

     if st.button("🔄 Retake Quiz"):
         st.session_state.user_answers = {}
         st.rerun()

# Footer
st.markdown("---")
st.caption("Powered by Hugging Face & Smart Text Analysis | Supports PDF, DOCX, PPTX")

if __name__ == "__main__":
    main()