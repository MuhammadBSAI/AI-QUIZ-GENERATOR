import fitz  # PyMuPDF
import docx
from pptx import Presentation


def read_pdf(file):
    text = ""
    pdf = fitz.open(stream=file.read(), filetype="pdf")
    for page in pdf:
        text += page.get_text()
    return text


def read_docx(file):
    doc = docx.Document(file)
    return "\n".join([p.text for p in doc.paragraphs])


def read_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text